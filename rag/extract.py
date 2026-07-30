"""Extract text from PDF, DOCX, XLSX, and PPTX files for RAG indexing."""

from __future__ import annotations

import logging
from pathlib import Path

log = logging.getLogger(__name__)


def extract_pdf(data: bytes) -> str:
    import pymupdf
    doc = pymupdf.open(stream=data, filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n\n".join(pages)


def extract_docx(data: bytes) -> str:
    import io
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def extract_xlsx(data: bytes) -> str:
    import io
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"## Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    wb.close()
    return "\n\n".join(parts)


def extract_pptx(data: bytes) -> str:
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    slide_text.append(" | ".join(cells))
        if slide_text:
            parts.append(f"## Slide {i}\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


_EXTRACTORS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "xlsx": extract_xlsx,
    "pptx": extract_pptx,
}


def extract_text(filename: str, data: bytes) -> str | None:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    extractor = _EXTRACTORS.get(ext)
    if not extractor:
        log.warning("no extractor for .%s — skipping %s", ext, filename)
        return None
    try:
        text = extractor(data)
        if not text or not text.strip():
            log.warning("no text extracted from %s", filename)
            return None
        return text.strip()
    except Exception:
        log.exception("failed to extract text from %s", filename)
        return None


def doc_type_from_filename(filename: str) -> str:
    name = filename.lower()
    if "cim" in name or "memorandum" in name:
        return "cim"
    if "term" in name and "sheet" in name:
        return "legal"
    if "pitch" in name or "deck" in name or "presentation" in name:
        return "cim"
    if "qoe" in name or "quality" in name:
        return "qoe"
    if "msa" in name or "agreement" in name or "contract" in name:
        return "msa"
    if "esg" in name:
        return "esg"
    if "tax" in name:
        return "tax"
    if "ddq" in name or "diligence" in name:
        return "tech_ddq"
    if "startup" in name or "database" in name:
        return "industry"
    return "misc"
