"""Generate a Lithuanian FastVC product tour PPTX from *_lt.png screenshots.

Output: docs/fastvc-product-tour-lt.pptx

Usage:
    python -m scripts.make_pptx_lt
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "screenshots"
OUT = ROOT / "docs" / "fastvc-product-tour-lt.pptx"

INK = RGBColor(0x14, 0x23, 0x1B)
INK_MUTED = RGBColor(0x41, 0x50, 0x46)
ACCENT = RGBColor(0x1F, 0x5D, 0x43)
BG = RGBColor(0xF7, 0xF6, 0xF1)

SLIDES = [
    {
        "eyebrow": "PEHERO",
        "title": "Jūsų privataus kapitalo\nAI agentų komanda",
        "subtitle": "Paieška · vertinimas · patikra · kapitalas · portfelio operacijos",
        "screenshot": None,
    },
    {
        "eyebrow": "01 · POKALBIS",
        "title": "Vienas pokalbis, kiekvienas VC specialistas",
        "subtitle": "Įveskite prefiksą arba klauskite lietuviškai — FastVC parinks tinkamą agentą.",
        "screenshot": "04-chat-empty_lt.png",
    },
    {
        "eyebrow": "01 · POKALBIS",
        "title": "Sandorių atranka — taip / ne per 90 sekundžių",
        "subtitle": "DR VET — tikra lietuviška veterinarijos klinika su tikrais finansiniais duomenimis.",
        "screenshot": "05-chat-triage_lt.png",
    },
    {
        "eyebrow": "01 · POKALBIS",
        "title": "LTM finansinių rodiklių normalizavimas",
        "subtitle": "Tikri DR VET finansiniai duomenys iš duomenų bazės — pajamos, EBITDA, augimas.",
        "screenshot": "06-chat-ltm_lt.png",
    },
    {
        "eyebrow": "01 · POKALBIS",
        "title": "IC Memo rengėjas",
        "subtitle": "Investicijų komiteto memo, parengtas iš Kardiolita sandorio duomenų.",
        "screenshot": "07-chat-memo_lt.png",
    },
    {
        "eyebrow": "02 · PIPELINE",
        "title": "Kanban per visus sandorių etapus",
        "subtitle": "157 tikrų Lietuvos įmonių — sveikatos, draudimo, logistikos, NT sektoriuose.",
        "screenshot": "08-pipeline-kanban_lt.png",
    },
    {
        "eyebrow": "03 · SANDORIO DETALĖS",
        "title": "Kiekvienas sandoris turi savo darbo erdvę",
        "subtitle": "Aprašymas dešinėje, pokalbis centre, artefaktai srautu ateina.",
        "screenshot": "09-pipeline-deal_lt.png",
    },
    {
        "eyebrow": "04 · ANALITIKA",
        "title": "Klauskite lietuviškai, gaukite grafiką",
        "subtitle": "Top 10 įmonių pagal pajamas — tikri duomenys iš rekvizitai.vz.lt.",
        "screenshot": "11-analytics-revenue_lt.png",
    },
    {
        "eyebrow": "05 · INSTRUKCIJOS",
        "title": "Koreguokite komandą gyvai",
        "subtitle": "Kiekvieno specialisto instrukcijos redaguojamos iš tos pačios sąsajos.",
        "screenshot": "12-instructions-list_lt.png",
    },
    {
        "eyebrow": "SUSISIEKITE",
        "title": "Pamatykite FastVC\nsu savo sandoriais",
        "subtitle": "hello@fastvc.fyi · fastvc.fyi/contact\nBYOD — naudokite savo sandorių duomenis.",
        "screenshot": None,
    },
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for s in SLIDES:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG

        if s["screenshot"] and (SHOTS / s["screenshot"]).exists():
            # Two-column: text left, screenshot right
            left = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(4.5), Inches(5.5))
            tf = left.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = s["eyebrow"]
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.space_after = Pt(8)

            p = tf.add_paragraph()
            p.text = s["title"]
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = INK
            p.space_after = Pt(12)

            p = tf.add_paragraph()
            p.text = s["subtitle"]
            p.font.size = Pt(13)
            p.font.color.rgb = INK_MUTED
            p.space_after = Pt(6)

            img_path = str(SHOTS / s["screenshot"])
            slide.shapes.add_picture(img_path, Inches(5.3), Inches(0.5), Inches(7.7))
        else:
            # Full-width centered (hero / closing)
            tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5))
            tf = tb.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = s["eyebrow"]
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(12)

            p = tf.add_paragraph()
            p.text = s["title"]
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = INK
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(16)

            p = tf.add_paragraph()
            p.text = s["subtitle"]
            p.font.size = Pt(16)
            p.font.color.rgb = INK_MUTED
            p.alignment = PP_ALIGN.CENTER

        # Footer
        footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12), Inches(0.4))
        p = footer.text_frame.paragraphs[0]
        p.text = "FastVC · Jūsų privataus kapitalo AI agentų komanda"
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(0x7A, 0x86, 0x7E)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}  ({OUT.stat().st_size / 1024:.1f} KB)")


if __name__ == "__main__":
    build()
